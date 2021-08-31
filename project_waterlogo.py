from wand.image import Image as I
from PIL import Image
from pptx import Presentation
def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]
    im = Image.open(image_url)
    width, height = im.size
    placeholder.height = height
    placeholder.width = width
    placeholder = placeholder.insert_picture(image_url)
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
for i in range(1,6):
    with I(filename=f"d:/PPT - Logo - Images - Assignment/image{i}.jpg") as image:
        with I(filename ='d:/PPT - Logo - Images - Assignment/nike_black.png') as water:
            with image.clone() as watermark:
                watermark.watermark(water,0,10, 20)
                watermark.save(filename =f'd:/PPT - Logo - Images - Assignment/logo_added{i}.jpg')
prs=Presentation()
for i in range(1,6):
    lyt=prs.slide_layouts[8] 
    slide=prs.slides.add_slide(lyt)
    title=slide.shapes.title 
    subtitle=slide.placeholders[2] 
    title.text="Assigment" 
    subtitle.text=f"Image_watermark{i}"
    img_path=f"d:/PPT - Logo - Images - Assignment/logo_added{i}.jpg"
    _add_image(slide,1,img_path)
    prs.save("d:/PPT - Logo - Images - Assignment/slide1.pptx")  